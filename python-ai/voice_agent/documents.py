import json
import os
import random
import re

import requests
from openpyxl import Workbook
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

OLLAMA_URL = "http://localhost:11434/api/chat"
MODEL = "llama3.2"

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "generated_docs")

FONT_REGULAR = "C:\\Windows\\Fonts\\arial.ttf"
FONT_BOLD = "C:\\Windows\\Fonts\\arialbd.ttf"


def _ask_ollama(system_prompt, user_text, as_json=True):
    payload = {
        "model": MODEL,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_text},
        ],
        "stream": False,
        "options": {"temperature": 0.4},
    }
    if as_json:
        payload["format"] = "json"
    resp = requests.post(OLLAMA_URL, json=payload, timeout=180)
    resp.raise_for_status()
    return resp.json()["message"]["content"]


def slugify(text):
    slug = re.sub(r"[^a-z0-9]+", "-", (text or "").lower()).strip("-")
    return slug or "document"


# ---------- Excel ----------

EXCEL_PROMPT = """Generate simple tabular data for the given topic.
Output ONLY a JSON object: {"headers": ["col1", "col2", ...], "rows": [["a","b",...], ...]}
Keep it to 5-15 data rows and 2-6 columns. No explanations."""


def create_excel(topic):
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    raw = _ask_ollama(EXCEL_PROMPT, topic, as_json=True)
    data = json.loads(raw)

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(data.get("headers", []))
    for row in data.get("rows", []):
        ws.append(row)

    path = os.path.join(OUTPUT_DIR, f"{slugify(topic)}.xlsx")
    wb.save(path)
    return path


# ---------- Word ----------

WORD_PROMPT = """Write short, well-organized content for the given topic, suitable for a document/post.
Output plain text only: a short title on the first line, then 2-5 paragraphs. No markdown, no JSON."""


def create_word(topic):
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    raw = _ask_ollama(WORD_PROMPT, topic, as_json=False)
    lines = [l.strip() for l in raw.strip().split("\n") if l.strip()]
    title = lines[0] if lines else topic
    paragraphs = lines[1:] if len(lines) > 1 else []

    doc = Document()
    doc.add_heading(title, level=1)
    for para in paragraphs:
        doc.add_paragraph(para)

    path = os.path.join(OUTPUT_DIR, f"{slugify(topic)}.docx")
    doc.save(path)
    return path


# ---------- PowerPoint ----------

PPTX_PROMPT = """Create presentation slide content for the given topic.
Output ONLY a JSON object: {"title": "...", "slides": [{"title": "...", "bullets": ["...", "..."]}]}
4-8 slides, 2-5 bullets each. No explanations."""


def create_pptx(topic):
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    raw = _ask_ollama(PPTX_PROMPT, topic, as_json=True)
    data = json.loads(raw)

    prs = Presentation()

    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = data.get("title", topic)

    bullet_layout = prs.slide_layouts[1]
    for slide_data in data.get("slides", []):
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = slide_data.get("title", "")
        body = slide.placeholders[1].text_frame
        bullets = slide_data.get("bullets", [])
        if bullets:
            body.text = bullets[0]
            for bullet in bullets[1:]:
                p = body.add_paragraph()
                p.text = bullet

    path = os.path.join(OUTPUT_DIR, f"{slugify(topic)}.pptx")
    prs.save(path)
    return path


# ---------- Banner image ----------

PALETTES = [
    (("#FF6B6B", "#FFD93D"), "#FFFFFF"),
    (("#4FD6FF", "#6C5CE7"), "#FFFFFF"),
    (("#00C9A7", "#845EC2"), "#FFFFFF"),
    (("#F9844A", "#F72585"), "#FFFFFF"),
]


def _gradient(size, color_a, color_b):
    img = Image.new("RGB", size)
    draw = ImageDraw.Draw(img)
    ra, ga, ba = tuple(int(color_a[i:i + 2], 16) for i in (1, 3, 5))
    rb, gb, bb = tuple(int(color_b[i:i + 2], 16) for i in (1, 3, 5))
    width, height = size
    for x in range(width):
        t = x / max(width - 1, 1)
        r = int(ra + (rb - ra) * t)
        g = int(ga + (gb - ga) * t)
        b = int(ba + (bb - ba) * t)
        draw.line([(x, 0), (x, height)], fill=(r, g, b))
    return img


def _add_confetti(draw, size, count=40):
    for _ in range(count):
        x = random.randint(0, size[0])
        y = random.randint(0, size[1])
        r = random.randint(4, 10)
        color = random.choice(["#FFFFFF", "#FFD93D", "#FFFFFF"])
        draw.ellipse((x - r, y - r, x + r, y + r), fill=color)


def _draw_centered_text(draw, text, size, text_color, max_font_size=90, min_font_size=30):
    font_size = max_font_size
    font = ImageFont.truetype(FONT_BOLD, font_size)
    while True:
        bbox = draw.textbbox((0, 0), text, font=font)
        w, h = bbox[2] - bbox[0], bbox[3] - bbox[1]
        if w <= size[0] * 0.85 or font_size <= min_font_size:
            break
        font_size -= 6
        font = ImageFont.truetype(FONT_BOLD, font_size)

    x = (size[0] - w) / 2 - bbox[0]
    y = (size[1] - h) / 2 - bbox[1]
    draw.text((x, y), text, font=font, fill=text_color, stroke_width=3, stroke_fill="#00000055")


def create_banner(text, size=(1200, 600)):
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    (color_a, color_b), text_color = random.choice(PALETTES)
    img = _gradient(size, color_a, color_b)
    draw = ImageDraw.Draw(img)
    _add_confetti(draw, size)
    _draw_centered_text(draw, text, size, text_color)

    path = os.path.join(OUTPUT_DIR, f"{slugify(text)}-banner.png")
    img.save(path)
    return path


# ---------- Greeting / celebration card (PowerPoint) ----------

CARD_MESSAGE_PROMPT = """Write one short, warm greeting-card message (1-2 sentences, max 25 words)
for the given occasion. Plain text only, no quotes, no explanations."""


def create_card_pptx(topic):
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    message = _ask_ollama(CARD_MESSAGE_PROMPT, topic, as_json=False).strip().strip('"')

    slide_px = (1280, 720)

    (color_a1, color_b1), text_color1 = random.choice(PALETTES)
    cover = _gradient(slide_px, color_a1, color_b1)
    cover_draw = ImageDraw.Draw(cover)
    _add_confetti(cover_draw, slide_px, count=60)
    _draw_centered_text(cover_draw, topic, slide_px, text_color1)
    cover_path = os.path.join(OUTPUT_DIR, "_tmp_card_cover.png")
    cover.save(cover_path)

    (color_a2, color_b2), _ = random.choice([p for p in PALETTES if p[0] != (color_a1, color_b1)])
    inside_bg = _gradient(slide_px, color_a2, color_b2)
    inside_draw = ImageDraw.Draw(inside_bg)
    _add_confetti(inside_draw, slide_px, count=30)
    inside_path = os.path.join(OUTPUT_DIR, "_tmp_card_inside.png")
    inside_bg.save(inside_path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    slide1 = prs.slides.add_slide(blank_layout)
    slide1.shapes.add_picture(cover_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    slide2 = prs.slides.add_slide(blank_layout)
    slide2.shapes.add_picture(inside_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    plaque = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(2.4), prs.slide_width - Inches(3), Inches(2.7),
    )
    plaque.fill.solid()
    plaque.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    plaque.line.fill.background()
    plaque.shadow.inherit = False

    tf = plaque.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = message
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    path = os.path.join(OUTPUT_DIR, f"{slugify(topic)}-card.pptx")
    prs.save(path)

    os.remove(cover_path)
    os.remove(inside_path)
    return path
